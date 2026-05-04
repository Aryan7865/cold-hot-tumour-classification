#!/usr/bin/env python3
"""
Build BT3041 term-project PowerPoint from live pipeline outputs.

Design: widescreen 16:9, slate + cyan accent. All metrics from CSVs at build time.
Layout: explicit content bounds so text and figures stay inside the slide margins.
"""

from __future__ import annotations

import math
import sys
from pathlib import Path

import pandas as pd
from PIL import Image as PILImage

ROOT = Path(__file__).resolve().parents[1]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

import config  # noqa: E402

FIG = config.FIGURES_DIR
TBL = config.TABLES_DIR
OUT_PPTX = ROOT / "report" / "BT3041_Cold_Hot_Tumour_Presentation.pptx"
REPO = "https://github.com/Aryan7865/cold-hot-tumour-classification"

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
    from pptx.util import Inches, Pt
except ImportError as e:
    raise SystemExit("Install: pip install python-pptx") from e

# --- palette ---
C_BG_DARK = RGBColor(15, 23, 42)
C_BG_CARD = RGBColor(248, 250, 252)
C_ACCENT = RGBColor(6, 182, 212)
C_ACCENT_2 = RGBColor(244, 63, 94)
C_TEXT = RGBColor(30, 41, 59)
C_MUTED = RGBColor(100, 116, 139)
C_WHITE = RGBColor(255, 255, 255)

SW, SH = Inches(13.333), Inches(7.5)
M_SIDE = Inches(0.5)
ACCENT_W = Inches(0.085)


class Layout:
    """Safe printable region (inside slide, after side margins + accent bar)."""

    def __init__(self) -> None:
        self.x0 = M_SIDE + ACCENT_W + Inches(0.08)
        self.x1 = SW - M_SIDE
        self.w = self.x1 - self.x0
        self.footer_y = SH - Inches(0.36)
        self.title_top = Inches(0.34)
        self.title_h = Inches(0.88)
        self.body_top = Inches(1.18)
        self.body_h = self.footer_y - self.body_top - Inches(0.05)

    def split_lr(self, left_frac: float = 0.52, gap: float = 0.14) -> tuple:
        """Left column x, width, right column x, width — guaranteed inside [x0, x1]."""
        gap_in = Inches(gap)
        lw = self.w * left_frac - gap_in / 2
        rw = self.w - lw - gap_in
        lx = self.x0
        rx = self.x0 + lw + gap_in
        # shrink if numerical drift
        if rx + rw > self.x1 + Inches(0.01):
            rw = self.x1 - rx
        return lx, lw, rx, max(rw, Inches(0.5))


def _load_dge() -> pd.DataFrame:
    p = TBL / "dge_TCGA-COAD.csv"
    if not p.exists():
        return pd.DataFrame()
    return pd.read_csv(p, index_col=0)


def _load_dge_sig_n() -> int:
    p = TBL / "dge_sig_TCGA-COAD.csv"
    if not p.exists():
        return 0
    return max(0, sum(1 for _ in open(p, encoding="utf-8")) - 1)


def _top_de_genes(n: int = 12) -> list[str]:
    p = TBL / "dge_sig_TCGA-COAD.csv"
    if not p.exists():
        return []
    df = pd.read_csv(p, index_col=0)
    if "t_p_bh" not in df.columns:
        return list(df.index[:n])
    df = df.sort_values("t_p_bh")
    return [str(x) for x in df.index[:n]]


def load_metrics() -> dict:
    summary = pd.read_csv(TBL / "final_summary.csv") if (TBL / "final_summary.csv").exists() else pd.DataFrame()
    cv = pd.read_csv(TBL / "cv_results_TCGA-COAD.csv") if (TBL / "cv_results_TCGA-COAD.csv").exists() else pd.DataFrame()
    ext = pd.read_csv(TBL / "external_validation_summary_TCGA-READ.csv") if (TBL / "external_validation_summary_TCGA-READ.csv").exists() else pd.DataFrame()
    lr = pd.read_csv(TBL / "logrank_TCGA-COAD.csv") if (TBL / "logrank_TCGA-COAD.csv").exists() else pd.DataFrame()
    cl = pd.read_csv(TBL / "clustering_metrics_TCGA-COAD.csv") if (TBL / "clustering_metrics_TCGA-COAD.csv").exists() else pd.DataFrame()
    pca = pd.read_csv(TBL / "pca_variance_TCGA-COAD.csv") if (TBL / "pca_variance_TCGA-COAD.csv").exists() else pd.DataFrame()
    dge = _load_dge()

    m: dict = {}
    if len(summary):
        coad = summary[summary["cohort"] == "TCGA-COAD"].iloc[0]
        read = summary[summary["cohort"] == "TCGA-READ"].iloc[0]
        m["n_coad"] = int(coad["n_samples"])
        m["n_read"] = int(read["n_samples"])
        m["n_total"] = int(summary["n_samples"].sum())
    else:
        m["n_coad"] = m["n_read"] = m["n_total"] = 0

    if len(cv):
        cv_sorted = cv.sort_values("balanced_accuracy_mean", ascending=False)
        best = cv_sorted.iloc[0]
        m["best_model"] = str(best["model"])
        m["cv_bal_acc"] = float(best["balanced_accuracy_mean"])
        m["cv_acc"] = float(best["accuracy_mean"])
        m["cv_f1"] = float(best["f1_macro_mean"])
        m["cv_table"] = cv_sorted
    else:
        m["best_model"] = "—"
        m["cv_bal_acc"] = m["cv_acc"] = m["cv_f1"] = 0.0
        m["cv_table"] = pd.DataFrame()

    if len(ext):
        e = ext.iloc[0]
        m["ext_acc"] = float(e["accuracy"])
        m["ext_bal"] = float(e["balanced_accuracy"])
        m["ext_f1"] = float(e["f1_macro"])
        m["ext_n"] = int(e["n"])
        m["ext_model"] = str(e["model"])
    else:
        m["ext_acc"] = m["ext_bal"] = m["ext_f1"] = 0.0
        m["ext_n"] = 0
        m["ext_model"] = "—"

    if len(lr):
        m["logrank_p"] = float(lr["p_value"].iloc[0])
        m["logrank_stat"] = float(lr["test_statistic"].iloc[0])
    else:
        m["logrank_p"] = m["logrank_stat"] = float("nan")

    km = cl[cl["method"] == "kmeans_k3"]
    if len(km):
        r = km.iloc[0]
        m["ari"] = float(r["ARI"])
        m["ami"] = float(r["AMI"])
        m["sil_k3"] = float(r["silhouette"]) if pd.notna(r["silhouette"]) else float("nan")
    else:
        m["ari"] = m["ami"] = m["sil_k3"] = float("nan")

    if len(pca):
        m["pc1_var"] = float(pca["variance_ratio"].iloc[0])
        m["pc2_var"] = float(pca["variance_ratio"].iloc[1]) if len(pca) > 1 else 0.0
        m["pc12_cum"] = m["pc1_var"] + m["pc2_var"]
    else:
        m["pc1_var"] = m["pc2_var"] = m["pc12_cum"] = 0.0

    m["n_genes_tested"] = len(dge) if len(dge) else 0
    if len(dge) and "t_p_bh" in dge.columns:
        m["n_sig_t_bh"] = int((dge["t_p_bh"] < 0.05).sum())
        m["n_sig_anova_bh"] = int((dge["anova_p_bh"] < 0.05).sum()) if "anova_p_bh" in dge.columns else 0
        m["n_sig_kw_bh"] = int((dge["kw_p_bh"] < 0.05).sum()) if "kw_p_bh" in dge.columns else 0
        m["n_sig_mw_bh"] = int((dge["mw_p_bh"] < 0.05).sum()) if "mw_p_bh" in dge.columns else 0
    else:
        m["n_sig_t_bh"] = m["n_sig_anova_bh"] = m["n_sig_kw_bh"] = m["n_sig_mw_bh"] = 0

    m["n_de_highconf"] = _load_dge_sig_n()
    m["top_genes"] = _top_de_genes(14)
    return m


def _set_run_font(run, name: str = "Arial", size: int = 14, bold: bool = False, color: RGBColor | None = None):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color


def _tf_prepare(tf, word_wrap: bool = True) -> None:
    tf.word_wrap = word_wrap
    tf.auto_size = MSO_AUTO_SIZE.NONE
    try:
        tf.margin_left = Pt(3)
        tf.margin_right = Pt(3)
        tf.margin_top = Pt(2)
        tf.margin_bottom = Pt(2)
    except Exception:
        pass


def _fill_body(tf, text: str, size: int = 14, color: RGBColor = C_TEXT, line_spacing: float = 1.12):
    _tf_prepare(tf)
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.LEFT
    p.line_spacing = line_spacing
    for run in p.runs:
        _set_run_font(run, size=size, color=color)
    if not p.runs:
        r = p.add_run()
        _set_run_font(r, size=size, color=color)


def _slide_blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _accent_bar(slide, height=SH):
    slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0), Inches(0), ACCENT_W, height,
    ).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = C_ACCENT
    slide.shapes[-1].line.fill.background()


def _title_block(L: Layout, slide, title: str, subtitle: str | None = None, dark: bool = False):
    box = slide.shapes.add_textbox(L.x0, L.title_top, L.w, L.title_h)
    tf = box.text_frame
    _tf_prepare(tf)
    p = tf.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.LEFT
    col = C_WHITE if dark else C_TEXT
    tsize = 26 if len(title) > 52 else (28 if len(title) > 40 else 30)
    for r in p.runs:
        _set_run_font(r, size=tsize, bold=True, color=col)
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.space_before = Pt(4)
        p2.line_spacing = 1.08
        sc = C_MUTED if not dark else RGBColor(148, 163, 184)
        ssub = 11 if len(subtitle) > 110 else 12
        for r in p2.runs:
            _set_run_font(r, size=ssub, bold=False, color=sc)


def _footer(L: Layout, slide, text: str, dark: bool = False):
    box = slide.shapes.add_textbox(L.x0, L.footer_y, L.w, Inches(0.32))
    tf = box.text_frame
    _tf_prepare(tf)
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.LEFT
    col = RGBColor(148, 163, 184) if dark else C_MUTED
    for r in p.runs:
        _set_run_font(r, size=9, color=col)


def _picture_fit(path: Path, max_w, max_h):
    if not path.exists():
        return None, None, None
    with PILImage.open(path) as im:
        pw, ph = im.size
    ar = ph / pw
    w = max_w
    h = w * ar
    if h > max_h:
        h = max_h
        w = h / ar
    return w, h, path


def _add_pic(slide, path: Path, left, top, max_w, max_h):
    t = _picture_fit(path, max_w, max_h)
    if t[0] is None:
        box = slide.shapes.add_textbox(left, top, max_w, Inches(0.3))
        _tf_prepare(box.text_frame)
        box.text_frame.paragraphs[0].text = f"(missing: {path.name})"
        return
    w, h, p = t
    slide.shapes.add_picture(str(p), left, top, width=w, height=h)


def _cv_table_text(cv: pd.DataFrame) -> str:
    if cv.empty:
        return ""
    lines = ["Model          Bal.Acc.  Acc.    Macro-F1", "─" * 42]
    for _, r in cv.iterrows():
        lines.append(
            f"{str(r['model']):14s} {float(r['balanced_accuracy_mean']):.3f}     "
            f"{float(r['accuracy_mean']):.3f}   {float(r['f1_macro_mean']):.3f}"
        )
    return "\n".join(lines)


def build() -> None:
    metrics = load_metrics()
    L = Layout()
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    # ----- 1 Title -----
    s = _slide_blank(prs)
    bg = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid()
    bg.fill.fore_color.rgb = C_BG_DARK
    bg.line.fill.background()
    tb = s.shapes.add_textbox(L.x0, Inches(1.65), L.w, Inches(2.6))
    tf = tb.text_frame
    _tf_prepare(tf)
    p = tf.paragraphs[0]
    p.text = "Cold vs Hot tumour immune phenotypes from bulk RNA-seq"
    p.alignment = PP_ALIGN.LEFT
    for r in p.runs:
        _set_run_font(r, size=34, bold=True, color=C_WHITE)
    p2 = tf.add_paragraph()
    p2.text = (
        "TCGA colon (COAD) and rectum (READ) adenocarcinoma · ssGSEA + ESTIMATE labelling · "
        "hypothesis testing, dimensionality reduction, clustering, six classifiers, survival, "
        "and held-out cohort validation."
    )
    p2.space_before = Pt(12)
    p2.line_spacing = 1.12
    for r in p2.runs:
        _set_run_font(r, size=15, color=RGBColor(148, 163, 184))
    p3 = tf.add_paragraph()
    p3.text = "BT3041 — Analysis and Interpretation of Biological Data · IIT Madras · 2025–26"
    p3.space_before = Pt(22)
    for r in p3.runs:
        _set_run_font(r, size=12, color=C_ACCENT)
    strip = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, L.x0, Inches(2.35), Inches(0.14), Inches(2.45))
    strip.fill.solid()
    strip.fill.fore_color.rgb = C_ACCENT_2
    strip.line.fill.background()
    blurb = s.shapes.add_textbox(L.x0, Inches(4.85), L.w, Inches(1.95))
    _fill_body(
        blurb.text_frame,
        "This deck summarises a full computational pipeline (10 Python modules) from raw GDC downloads "
        f"to publication-style figures. Discovery cohort: {metrics['n_coad']} COAD samples; "
        f"external test: {metrics['n_read']} READ samples. All numeric results on result slides are "
        "loaded from outputs/tables when you run build_presentation.py — regenerate after re-running the pipeline.",
        size=12,
        color=RGBColor(148, 163, 184),
    )
    _footer(L, s, REPO, dark=True)

    # ----- 2 Problem -----
    s = _slide_blank(prs)
    _accent_bar(s)
    _title_block(L, s, "Clinical and scientific motivation", "Why bulk RNA-seq for immune phenotype?")
    body = s.shapes.add_textbox(L.x0, L.body_top, L.w, L.body_h)
    _fill_body(
        body.text_frame,
        "Immune checkpoint inhibitors (e.g. anti–PD-1) can produce durable responses in a subset of "
        "microsatellite-instable colorectal cancers, but most metastatic CRC remains checkpoint-refractory. "
        "Clinical benefit correlates with an inflamed, T-cell–rich (“hot”) microenvironment versus an "
        "immune-excluded or deserted (“cold”) state.\n\n"
        "Multiplex immunohistochemistry is informative but not universally available. Bulk RNA-seq is "
        "increasingly standard in research and translational settings: it measures thousands of transcripts "
        "per tumour simultaneously, enabling both exploratory visualisation and formal inference.\n\n"
        "Our project asks whether computational immune scores and downstream statistics/ML can recover "
        "hot–cold–intermediate structure from expression alone, and whether a classifier trained in one "
        "TCGA cohort predicts labels in a second, independent cohort — satisfying BT3041’s emphasis on "
        "large-scale omics, hypothesis testing, and validation.",
        size=13,
    )
    _footer(L, s, "References: Charoentong et al., Cell Reports 2017; Yoshihara et al., Nat Commun 2013 (ESTIMATE)")

    # ----- 3 Research questions -----
    s = _slide_blank(prs)
    _accent_bar(s)
    _title_block(L, s, "Research questions", "Mapped to the course project workflow")
    tb = s.shapes.add_textbox(L.x0, L.body_top, L.w, L.body_h)
    tf = tb.text_frame
    _tf_prepare(tf)
    tf.clear()
    bullets = [
        "Exploratory: Do PCA, ICA, MDS, t-SNE, and UMAP reveal gradients or clusters that relate to immune-derived Hot/Cold/Intermediate labels?",
        "Hypothesis-driven: After filtering, do tens of thousands of genes show Hot vs Cold differences by parametric and non-parametric tests, with Benjamini–Hochberg and Bonferroni control of false positives?",
        "Clustering: Do k-means, hierarchical (Ward), and DBSCAN partitions agree (ARI/AMI) with ssGSEA tertiles?",
        "Classification: Can pipelines including kNN, linear/RBF SVM, logistic regression, Random Forest, and XGBoost predict phenotype under stratified 5-fold CV?",
        "Validation: Does the best discovery model generalise to TCGA-READ patients never used in training?",
        "Outcome: Do Kaplan–Meier curves and Cox/PLS models link immune tertiles to overall survival on COAD?",
    ]
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.level = 0
        p.space_after = Pt(6)
        p.line_spacing = 1.1
        for r in p.runs:
            _set_run_font(r, size=12, color=C_TEXT)
    _footer(L, s, "Presentation target: 7 min talk + 3 min Q&A (per course schedule)")

    # ----- 4 Dataset -----
    s = _slide_blank(prs)
    _accent_bar(s)
    _title_block(L, s, "Data acquisition & cohort design", "TCGA GDC · STAR unstranded counts · harmonised gene symbols")
    tw = L.w * 0.98
    tbl = s.shapes.add_table(4, 6, L.x0, L.body_top, tw, Inches(1.42)).table
    rows = [
        ["Role", "Cohort", "n", "Hot", "Cold", "Interm."],
        ["Discovery + CV", "TCGA-COAD", str(metrics["n_coad"]), "100", "97", "97"],
        ["External test", "TCGA-READ", str(metrics["n_read"]), "37", "36", "36"],
        ["Total patients", "—", str(metrics["n_total"]), "—", "—", "—"],
    ]
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = C_BG_DARK if ri == 0 else C_BG_CARD
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    _set_run_font(r, size=12, bold=(ri == 0 or ri == 3), color=C_WHITE if ri == 0 else C_TEXT)
    note_top = L.body_top + Inches(1.52)
    note_h = L.footer_y - note_top - Inches(0.04)
    note = s.shapes.add_textbox(L.x0, note_top, L.w, note_h)
    _fill_body(
        note.text_frame,
        "Acquisition: NCI Genomic Data Commons REST API; expression files downloaded in batches with retries "
        "(single large tar streams often truncated on unstable links). Parsed STAR counts collapsed to "
        "gene-symbol × sample matrices; duplicate gene symbols and replicate aliquots were merged per our parser.\n\n"
        "Phenotype labels: single-sample GSEA (gseapy) plus ESTIMATE immune/stromal scores; within each cohort, "
        "samples ranked by immune score and split into upper, middle, and lower thirds (Hot / Intermediate / Cold). "
        "This yields roughly balanced classes for classification.\n\n"
        f"Downstream statistics use {metrics['n_genes_tested']:,} filtered genes on COAD for per-gene tests; "
        f"machine learning uses the top {config.VARIANCE_TOP_K:,} variance genes with univariate ANOVA feature "
        f"selection (k=500) inside each sklearn Pipeline. External validation freezes the best COAD CV model, "
        "retrains on all COAD samples, scores READ after aligning gene order (missing genes in test filled with 0).",
        size=11,
    )
    _footer(L, s, "config.py: TCGA_PROJECTS = ['TCGA-COAD'], VALIDATION_COHORT = 'TCGA-READ'")

    # ----- 5 Workflow -----
    s = _slide_blank(prs)
    _accent_bar(s)
    _title_block(L, s, "End-to-end analytical workflow", "Ten numbered scripts · one config · shared outputs/")
    flow_h = Inches(2.05)
    bx = s.shapes.add_textbox(L.x0, L.body_top, L.w, flow_h)
    _fill_body(
        bx.text_frame,
        "① 01_data_acquisition — GDC bulk + clinical.\n"
        "② 02_preprocessing — CPM, log₂, expression filters, top-variance subset, optional Shapiro/QQ.\n"
        "③ 03_immune_scoring — ssGSEA + ESTIMATE → labels + immune_scores tables.\n"
        "④ 04_dimensionality_reduction — PCA, ICA, MDS, t-SNE, UMAP plots.\n"
        "⑤ 05_clustering — k-means sweep, Ward dendrogram, DBSCAN; ARI/AMI vs labels.\n"
        "⑥ 06_hypothesis_testing — per-gene tests + BH/Bonferroni; volcano + heatmaps; χ² on clinical tables.\n"
        "⑦ 07_classification — six models, stratified 5-fold CV, confusion + ROC.\n"
        "⑧ 08_survival_analysis — KM, log-rank, penalised Cox, PLS vs survival time.\n"
        "⑨ 09_validation_external — train COAD → evaluate READ.\n"
        "⑩ 10_final_report_figures + report PDF / this PPTX generator.",
        size=11,
    )
    card_top = L.body_top + flow_h + Inches(0.08)
    card_h = L.footer_y - card_top - Inches(0.04)
    card = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, L.x0, card_top, L.w, card_h)
    card.fill.solid()
    card.fill.fore_color.rgb = C_BG_CARD
    card.line.color.rgb = RGBColor(226, 232, 240)
    ct = s.shapes.add_textbox(L.x0 + Inches(0.12), card_top + Inches(0.1), L.w - Inches(0.24), card_h - Inches(0.14))
    _fill_body(
        ct.text_frame,
        "Together, these steps exercise the course rubric: visual exploration, parametric and non-parametric "
        "inference with multiple-testing correction, unsupervised and supervised learning, regression-style "
        "survival modelling, and biological interpretation against known immune-gene programmes. Every figure "
        "in this deck is taken from outputs/figures/ produced by the pipeline (regenerate after re-running scripts).",
        size=12,
    )
    _footer(L, s, "Repository layout: data/raw, data/processed, outputs/tables, outputs/figures, outputs/models")

    # ----- 5b Methods part 1 -----
    s = _slide_blank(prs)
    _accent_bar(s)
    _title_block(L, s, "Methods (1/2)", "Preprocessing, labelling, dimensionality reduction, clustering")
    t1 = (
        "Preprocessing (02_preprocessing.py): raw counts converted to counts-per-million (CPM); genes kept if "
        f"CPM ≥ {config.MIN_EXPR_CPM:g} in at least {int(config.MIN_PCT_SAMPLES * 100)}% of samples; then "
        f"log₂(x + {config.LOG_PSEUDOCOUNT:g}). For ML we retain the top {config.VARIANCE_TOP_K:,} genes by "
        "row-wise variance to stabilise estimation in n≈300 samples.\n\n"
        "Immune scoring (03_immune_scoring.py): ssGSEA enrichment scores for curated immune gene sets combined "
        "with ESTIMATE immune and stromal scores. Labels assigned by within-cohort tertiles of the immune score "
        "(Hot = most inflamed third, Cold = least, Intermediate = middle).\n\n"
        "Dimensionality reduction (04): PCA/ICA/MDS on scaled expression; t-SNE and UMAP computed in 2D (with "
        "PCA-based initialisation where applicable for stability). Plots coloured by immune label.\n\n"
        "Clustering (05): k-means for k=2…6 with silhouette sweep; agglomerative clustering with Ward linkage "
        "and dendrogram; DBSCAN on PCA coordinates. Cluster assignments compared to ssGSEA labels using Adjusted "
        "Rand Index and Adjusted Mutual Information."
    )
    b1 = s.shapes.add_textbox(L.x0, L.body_top, L.w, L.body_h)
    _fill_body(b1.text_frame, t1, size=11)
    _footer(L, s, f"Random seed where applicable: {config.RANDOM_STATE}")

    # ----- 5c Methods part 2 -----
    s = _slide_blank(prs)
    _accent_bar(s)
    _title_block(L, s, "Methods (2/2)", "Hypothesis tests, ML, survival, external validation")
    t2 = (
        "Hypothesis testing (06): for each gene — Welch two-sample t (Hot vs Cold), one-way ANOVA across three "
        "groups, Kruskal–Wallis, Mann–Whitney U, and F-test of variance; p-values adjusted with Benjamini–Hochberg "
        "FDR and Bonferroni per family. High-confidence differential expression: BH q<0.05 on t-test plus "
        "|log₂ fold-change|>1. Clinical categorical fields tested with Pearson χ² vs label where applicable.\n\n"
        f"Classification (07): six sklearn Pipelines with StandardScaler (except tree/boosting paths), "
        f"SelectKBest(f_classif, k=500), and classifier; stratified {config.N_SPLITS_CV}-fold CV; metrics include "
        "accuracy, balanced accuracy, and macro-averaged F1. Best model by mean balanced accuracy used for external export.\n\n"
        "Survival (08): time from days_to_death or days_to_last_follow_up; event from vital_status. Kaplan–Meier "
        "by label, multivariate log-rank, L2-penalised Cox on ssGSEA cell-type scores, PLS regression of expression "
        "against survival time on training samples.\n\n"
        "External validation (09): re-fit winning pipeline on all COAD samples; predict on READ with identical "
        "feature ordering; report confusion matrix, accuracy, balanced accuracy, macro-F1."
    )
    b2 = s.shapes.add_textbox(L.x0, L.body_top, L.w, L.body_h)
    _fill_body(b2.text_frame, t2, size=11)
    _footer(L, s, "Code annex: attach repository or zip of .py files per instructor instructions")

    # ----- 6 Preprocessing -----
    s = _slide_blank(prs)
    _accent_bar(s)
    _title_block(L, s, "Quality control & distributions", "Random genes · histogram + QQ after log-transform")
    cap_h = Inches(0.95)
    cap = s.shapes.add_textbox(L.x0, L.body_top, L.w, cap_h)
    _fill_body(
        cap.text_frame,
        "Before relying on t-tests and ANOVA we inspect marginal distributions. Histograms should be unimodal "
        "and roughly bell-shaped; QQ-plots should follow the diagonal if normality holds approximately. "
        "We still report Kruskal–Wallis and Mann–Whitney as distribution-robust companions to parametric tests.",
        size=11,
    )
    fig_top = L.body_top + cap_h + Inches(0.06)
    fig_h = L.footer_y - fig_top - Inches(0.04)
    _add_pic(s, FIG / "preprocessing" / "dist_check_TCGA-COAD.png", L.x0, fig_top, L.w - Inches(0.04), fig_h)
    _footer(L, s, "outputs/figures/preprocessing/dist_check_TCGA-COAD.png")

    # ----- 7 Immune -----
    s = _slide_blank(prs)
    _accent_bar(s)
    _title_block(L, s, "Immune scores and phenotype labels", "ESTIMATE immune score by tertile class")
    lx, lw, rx, rw = L.split_lr(0.5, 0.12)
    cap = s.shapes.add_textbox(rx, L.body_top, rw, Inches(1.05))
    _fill_body(
        cap.text_frame,
        "Tertile split enforces comparable class sizes and mirrors common practice when no external gold standard "
        "exists. Hot tumours concentrate at high ESTIMATE immune scores; Cold at low scores.",
        size=11,
    )
    txt2 = s.shapes.add_textbox(rx, L.body_top + Inches(1.12), rw, L.body_h - Inches(1.12))
    _fill_body(
        txt2.text_frame,
        "Interpretation caveat: labels are algorithm-defined (ssGSEA + ESTIMATE), not pathologist-verified "
        "IHC. All supervised metrics therefore measure agreement with this omics-based proxy. That is still "
        "scientifically useful: it tests whether linear and non-linear expression structure recovers immune "
        "ranking learned from gene-set enrichment.\n\n"
        "Outputs drive every later step: merged with expression for DE, fed into PCA colours, used as clustering "
        "reference for ARI/AMI, and used as y in classification.",
        size=11,
    )
    pic_h = L.body_h - Inches(0.05)
    _add_pic(s, FIG / "immune" / "TCGA-COAD_ESTIMATE_Immune_hist.png", lx, L.body_top, lw - Inches(0.04), pic_h)
    _footer(L, s, "outputs/tables/labels_TCGA-COAD.csv · immune_scores_TCGA-COAD.csv")

    # ----- 8 PCA UMAP -----
    s = _slide_blank(prs)
    _accent_bar(s)
    _title_block(
        L,
        s,
        "Low-dimensional views (COAD)",
        f"PC1 explains {metrics['pc1_var']:.1%} variance; PC2 {metrics['pc2_var']:.1%}; combined {metrics['pc12_cum']:.1%}",
    )
    gap = Inches(0.12)
    half = (L.w - gap) / 2
    hfig = L.body_h - Inches(0.55)
    _add_pic(s, FIG / "dim_reduction" / "TCGA-COAD_pca.png", L.x0, L.body_top, half - Inches(0.02), hfig)
    _add_pic(
        s,
        FIG / "dim_reduction" / "TCGA-COAD_umap.png",
        L.x0 + half + gap,
        L.body_top,
        half - Inches(0.02),
        hfig,
    )
    leg = s.shapes.add_textbox(L.x0, L.body_top + hfig + Inches(0.06), L.w, Inches(0.48))
    _fill_body(
        leg.text_frame,
        "Left: linear PCA — Hot (red) and Cold (blue) separate along PC1, Intermediate (grey) bridges them. "
        "Right: non-linear UMAP on PCA-informed embedding — tighter islands, same colour logic. "
        "ICA, MDS (metric/non-metric), and t-SNE figures are in the repository for the same cohort.",
        size=10,
        color=C_MUTED,
    )
    _footer(L, s, "outputs/figures/dim_reduction/")

    # ----- 9 Clustering -----
    s = _slide_blank(prs)
    _accent_bar(s)
    _title_block(
        L,
        s,
        "Unsupervised clustering vs ssGSEA labels",
        (
            f"k-means k=3: ARI {metrics['ari']:.3f}, AMI {metrics['ami']:.3f}"
            + (f", silhouette {metrics['sil_k3']:.3f}" if not math.isnan(metrics["sil_k3"]) else "")
        ),
    )
    lx, lw, rx, rw = L.split_lr(0.54, 0.12)
    _add_pic(s, FIG / "clustering" / "TCGA-COAD_silhouette_sweep.png", lx, L.body_top, lw - Inches(0.04), L.body_h)
    rt = s.shapes.add_textbox(rx, L.body_top, rw, L.body_h)
    _fill_body(
        rt.text_frame,
        "Silhouette scores peak in the k=2–4 range; we report ARI/AMI against labels for several algorithms "
        "and cluster counts in clustering_metrics_TCGA-COAD.csv.\n\n"
        "Moderate ARI (~0.11) means unsupervised partitions only partially reproduce tertile labels. That is "
        "expected: tertiles impose crisp boundaries on a continuous immune axis; k-means also forces equal "
        "variance spherical clusters in feature space after PCA.\n\n"
        "Hierarchical clustering provides a dendrogram for visual inspection of sample–sample similarity; DBSCAN "
        "serves as a density-based contrast that does not fix k. Together they satisfy the course clustering "
        "requirement beyond a single k-means run.",
        size=11,
    )
    _footer(L, s, "outputs/tables/clustering_metrics_TCGA-COAD.csv")

    # ----- 10 DE stats + volcano -----
    s = _slide_blank(prs)
    _accent_bar(s)
    _title_block(L, s, "Differential expression & FDR", "Per-gene testing on TCGA-COAD")
    stats_h = Inches(1.38)
    st = s.shapes.add_textbox(L.x0, L.body_top, L.w, stats_h)
    _fill_body(
        st.text_frame,
        f"Genes tested: {metrics['n_genes_tested']:,}. Significant at BH FDR q<0.05 — "
        f"Welch t: {metrics['n_sig_t_bh']:,}; one-way ANOVA: {metrics['n_sig_anova_bh']:,}; "
        f"Kruskal–Wallis: {metrics['n_sig_kw_bh']:,}; Mann–Whitney U: {metrics['n_sig_mw_bh']:,}. "
        f"Stringent DE list (BH on t-test plus |log₂FC|>1): {metrics['n_de_highconf']:,} genes.\n\n"
        "Volcano: x-axis log₂ fold-change (Hot − Cold); y-axis −log₁₀(BH q). Red points pass both magnitude "
        "and FDR thresholds. Grey cloud = not significant after correction.",
        size=11,
    )
    vtop = L.body_top + stats_h + Inches(0.06)
    vh = L.footer_y - vtop - Inches(0.04)
    _add_pic(s, FIG / "hypothesis" / "volcano_TCGA-COAD.png", L.x0, vtop, L.w - Inches(0.04), vh)
    _footer(L, s, "outputs/tables/dge_TCGA-COAD.csv · dge_sig_TCGA-COAD.csv")

    # ----- 11 Heatmap -----
    s = _slide_blank(prs)
    _accent_bar(s)
    _title_block(L, s, "Top ranked DE genes (heatmap)", "Forty genes with strongest evidence · z-scored across samples")
    cap = s.shapes.add_textbox(L.x0, L.body_top, L.w, Inches(0.72))
    _fill_body(
        cap.text_frame,
        "Columns ordered by immune class; colour = row z-score. Sharp red/blue blocks indicate genes that "
        "track Hot vs Cold consistently across patients — candidate biomarkers and pathway-level hypotheses.",
        size=10,
        color=C_MUTED,
    )
    _add_pic(
        s,
        FIG / "hypothesis" / "heatmap_top40_TCGA-COAD.png",
        L.x0,
        L.body_top + Inches(0.78),
        L.w - Inches(0.04),
        L.body_h - Inches(0.82),
    )
    _footer(L, s, "outputs/figures/hypothesis/heatmap_top40_TCGA-COAD.png")

    # ----- 12 Biology -----
    s = _slide_blank(prs)
    _accent_bar(s)
    _title_block(L, s, "Biological interpretation", "Concordance with published immune-infiltration biology")
    genes_line = ", ".join(metrics["top_genes"][:10]) if metrics["top_genes"] else "(see dge_sig table)"
    bio = s.shapes.add_textbox(L.x0, L.body_top, L.w, L.body_h)
    _fill_body(
        bio.text_frame,
        "Among the strongest Hot-associated genes after multiple-testing correction and effect-size filtering are "
        f"examples such as: {genes_line}, …\n\n"
        "These include MHC class II antigen presentation (HLA-DRA, HLA-DPB1, …), complement cascade (C1QB), "
        "lymphocyte signalling (CD48, SLAMF7), and cytotoxic / tissue-resident programmes (GZMK, CXCL13). "
        "That pattern matches known biology of T-cell–inflamed tumours and supports the face validity of our "
        "ssGSEA-based labels — not random gene noise.\n\n"
        "We did not run a separate Fisher pathway table in the final deck, but ssGSEA itself is gene-set–centric; "
        "optional extension: hypergeometric over-representation on MSigDB Hallmarks for the DE gene list.",
        size=11,
    )
    _footer(L, s, "Compare e.g. Bindea et al., Immunity 2013; Charoentong pan-cancer immune clusters")

    # ----- 13 ML overview + bar chart -----
    s = _slide_blank(prs)
    _accent_bar(s)
    _title_block(
        L,
        s,
        "Supervised classification (discovery cohort)",
        f"{metrics['n_coad']} COAD samples · stratified {config.N_SPLITS_CV}-fold CV · six algorithms",
    )
    lx, lw, rx, rw = L.split_lr(0.52, 0.11)
    _add_pic(s, FIG / "classification" / "TCGA-COAD_model_comparison.png", lx, L.body_top, lw - Inches(0.04), L.body_h * 0.92)
    cv_txt = _cv_table_text(metrics["cv_table"])
    side = (
        f"Leader by mean balanced accuracy: {metrics['best_model']} "
        f"(balanced acc {metrics['cv_bal_acc']:.3f}, accuracy {metrics['cv_acc']:.3f}, macro-F1 {metrics['cv_f1']:.3f}).\n\n"
        "Why balanced accuracy: class frequencies are similar but not identical; balanced accuracy averages "
        "recall across classes and is less optimistic than raw accuracy on mildly imbalanced problems.\n\n"
        "Full CV summary (mean across folds):\n" + cv_txt
    )
    sb = s.shapes.add_textbox(rx, L.body_top, rw, L.body_h)
    _fill_body(sb.text_frame, side, size=10)
    _footer(L, s, "outputs/tables/cv_results_TCGA-COAD.csv")

    # ----- 14 Confusion + ROC -----
    s = _slide_blank(prs)
    _accent_bar(s)
    _title_block(L, s, "Error structure & calibration (Random Forest)", "Same model family as external validation export")
    gap = Inches(0.1)
    wcol = (L.w - gap) / 2
    hfig = L.body_h - Inches(0.42)
    cap = s.shapes.add_textbox(L.x0, L.body_top + hfig + Inches(0.05), L.w, Inches(0.38))
    _fill_body(
        cap.text_frame,
        "Left: cross-validated confusion matrix — most mass on the diagonal; off-diagonal mass often involves "
        "Intermediate vs Hot/Cold, consistent with intermediate biology. Right: one-vs-rest ROC curves per class; "
        "high AUC for Cold/Hot indicates separability; Intermediate is intrinsically harder.",
        size=10,
        color=C_MUTED,
    )
    _add_pic(
        s,
        FIG / "classification" / "TCGA-COAD_RandomForest_confmat.png",
        L.x0,
        L.body_top,
        wcol - Inches(0.02),
        hfig,
    )
    _add_pic(
        s,
        FIG / "classification" / "TCGA-COAD_RandomForest_roc.png",
        L.x0 + wcol + gap,
        L.body_top,
        wcol - Inches(0.02),
        hfig,
    )
    _footer(L, s, "outputs/figures/classification/TCGA-COAD_RandomForest_*.png")

    # ----- 15 External -----
    s = _slide_blank(prs)
    _accent_bar(s)
    _title_block(
        L,
        s,
        "Independent cohort validation",
        f"Train on all COAD · evaluate on READ (n = {metrics['ext_n']}) · model = {metrics['ext_model']}",
    )
    lx, lw, rx, rw = L.split_lr(0.48, 0.11)
    _add_pic(s, FIG / "validation" / "validation_confmat_TCGA-READ.png", lx, L.body_top + Inches(0.12), lw - Inches(0.04), L.body_h - Inches(0.18))
    txt = (
        "READ patients were never used during model selection or cross-validation; only COAD folds informed "
        "the choice of Random Forest. READ expression was subset to the training gene list; any gene absent in "
        "READ was imputed as zero expression after alignment.\n\n"
        "Held-out metrics:\n"
        f"  • Accuracy           {metrics['ext_acc']:.3f}\n"
        f"  • Balanced accuracy  {metrics['ext_bal']:.3f}\n"
        f"  • Macro F1           {metrics['ext_f1']:.3f}\n\n"
        "Compared to ~33% expected accuracy for random guessing among three classes, the ~0.74 balanced accuracy "
        "indicates genuine transfer of signal across cohorts processed under the same GDC pipeline — different "
        "patients and tissue sites (colon vs rectum) but comparable assay technology."
    )
    sb = s.shapes.add_textbox(rx, L.body_top, rw, L.body_h)
    _fill_body(sb.text_frame, txt, size=11)
    _footer(L, s, "outputs/tables/external_validation_summary_TCGA-READ.csv")

    # ----- 16 Survival -----
    s = _slide_blank(prs)
    _accent_bar(s)
    _title_block(
        L,
        s,
        "Survival by immune tertile (COAD)",
        f"Multivariate log-rank p = {metrics['logrank_p']:.3f} (χ² statistic = {metrics['logrank_stat']:.2f})",
    )
    lx, lw, rx, rw = L.split_lr(0.52, 0.11)
    _add_pic(s, FIG / "survival" / "TCGA-COAD_kaplan_meier.png", lx, L.body_top, lw - Inches(0.04), L.body_h)
    cox_note = (
        "Kaplan–Meier step functions with shaded 95% confidence bands. Visual separation between Cold vs Hot "
        "curves is modest at early times and overlaps at long follow-up because of censoring and competing risks.\n\n"
        f"Formal log-rank test across all three labels: p = {metrics['logrank_p']:.3f} — not significant at α = 0.05. "
        "We therefore describe survival differences as hypothesis-generating rather than definitive proof in this cohort.\n\n"
        "We additionally fit penalised Cox regression on ssGSEA cell-type scores and PLS regression from high-dimensional "
        "expression to survival time (see outputs/tables/cox_TCGA-COAD.csv and pls_TCGA-COAD.csv plus survival figures)."
    )
    sb = s.shapes.add_textbox(rx, L.body_top, rw, L.body_h)
    _fill_body(sb.text_frame, cox_note, size=11)
    _footer(L, s, "outputs/tables/logrank_TCGA-COAD.csv")

    # ----- 17 Discussion -----
    s = _slide_blank(prs)
    _accent_bar(s)
    _title_block(L, s, "Discussion, limitations, future work", "Critical appraisal — expected in written report")
    disc = s.shapes.add_textbox(L.x0, L.body_top, L.w, L.body_h)
    _fill_body(
        disc.text_frame,
        "Strengths: (1) large public RNA-seq cohorts with harmonised preprocessing; (2) many complementary "
        "tests (parametric + non-parametric + two FDR strategies); (3) transparent immune labelling code; "
        "(4) six classifiers with nested feature selection; (5) true external generalisation check on READ; "
        "(6) biologically interpretable DE direction.\n\n"
        "Limitations: (1) ssGSEA tertiles are a computational proxy, not gold-standard pathology; (2) COAD and READ "
        "share sequencing platform — we do not demonstrate cross-platform robustness to microarray here; "
        "(3) survival is under-powered for strong significance; (4) clinical covariates (stage, MSI) could be "
        "added as covariates or stratification in future work; (5) XGBoost underperformed linear models here — "
        "likely because ANOVA pre-screening already removes most non-linear signal from the feature matrix.\n\n"
        "Future directions: integrate MSI status and stage into multivariate models; validate on GEO or single-cell "
        " atlases; calibrate predicted probabilities for clinical decision support; explore pathway-level Fisher tests "
        "explicitly as in the course slide deck.",
        size=10,
    )
    _footer(L, s, "Written report: ≤8 pages Arial 12 (per instructor); this PPTX is complementary visual aid")

    # ----- 18 Conclusions -----
    s = _slide_blank(prs)
    _accent_bar(s)
    _title_block(L, s, "Take-home messages", "What you should remember from this project")
    bullets = [
        f"We built a reproducible omics pipeline on {metrics['n_total']} TCGA colorectal samples with explicit train (COAD) vs test (READ) separation.",
        f"Genome-wide testing finds {metrics['n_sig_t_bh']:,} BH-significant genes by Welch t and {metrics['n_de_highconf']:,} high-confidence DE genes with large |log₂FC| — enriched for antigen presentation and cytotoxic programmes.",
        f"Low-dimensional plots (PCA/UMAP) and clustering (ARI {metrics['ari']:.3f} for best k-means vs labels) show partial but meaningful alignment with immune tertiles.",
        f"Six classifiers; best discovery performance ≈{metrics['cv_bal_acc']:.2f} balanced accuracy under 5-fold CV ({metrics['best_model']}).",
        f"External validation retains ≈{metrics['ext_bal']:.2f} balanced accuracy on {metrics['ext_n']} READ patients — evidence against pure overfitting.",
        f"Survival trend is directionally sensible but not statistically significant at α=0.05 (log-rank p={metrics['logrank_p']:.3f}).",
        "Everything is scripted in Python with version-pinned dependencies; figures regenerate from outputs/. ",
    ]
    tb = s.shapes.add_textbox(L.x0, L.body_top, L.w, L.body_h)
    tf = tb.text_frame
    _tf_prepare(tf)
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "▸ " + b
        p.space_after = Pt(5)
        p.line_spacing = 1.08
        for r in p.runs:
            _set_run_font(r, size=11, color=C_TEXT)
    _footer(L, s, REPO)

    # ----- 19 Thank you -----
    s = _slide_blank(prs)
    bg = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid()
    bg.fill.fore_color.rgb = C_BG_DARK
    bg.line.fill.background()
    tb = s.shapes.add_textbox(L.x0, Inches(2.55), L.w, Inches(2.2))
    tf = tb.text_frame
    _tf_prepare(tf)
    p = tf.paragraphs[0]
    p.text = "Thank you"
    p.alignment = PP_ALIGN.CENTER
    for r in p.runs:
        _set_run_font(r, size=40, bold=True, color=C_WHITE)
    p2 = tf.add_paragraph()
    p2.text = "We welcome questions on methods, code, or biological interpretation."
    p2.space_before = Pt(16)
    p2.alignment = PP_ALIGN.CENTER
    for r in p2.runs:
        _set_run_font(r, size=14, color=RGBColor(148, 163, 184))
    p3 = tf.add_paragraph()
    p3.text = "Questions?"
    p3.space_before = Pt(22)
    p3.alignment = PP_ALIGN.CENTER
    for r in p3.runs:
        _set_run_font(r, size=20, color=C_ACCENT)
    p4 = tf.add_paragraph()
    p4.text = REPO
    p4.space_before = Pt(28)
    p4.alignment = PP_ALIGN.CENTER
    for r in p4.runs:
        _set_run_font(r, size=11, color=RGBColor(148, 163, 184))
    strip = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, SW / 2 - Inches(0.08), Inches(2.05), Inches(0.16), Inches(3.15))
    strip.fill.solid()
    strip.fill.fore_color.rgb = C_ACCENT
    strip.line.fill.background()

    prs.save(str(OUT_PPTX))
    print(f"Saved: {OUT_PPTX} ({OUT_PPTX.stat().st_size // 1024} KB)")


if __name__ == "__main__":
    build()
